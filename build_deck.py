"""Re-runnable PowerPoint generator: how to publish a PWA to the Microsoft Store.

This deck documents the REUSABLE PROCESS for shipping any Progressive Web App
to the Microsoft Store - it is intentionally app-agnostic.

Run:
    C:/Users/vijayreddy/AppData/Local/Programs/Python/Python312/python.exe build_deck.py

Re-running overwrites Publishing-PWA-to-Microsoft-Store.pptx in place (idempotent).
To add a new step/milestone slide later, just append a dict to the SLIDES list
in the "MILESTONES" section below.
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Paths ------------------------------------------------------------------
PROJECT_ROOT = os.path.dirname(os.path.abspath(__file__))
OUTPUT_PATH = os.path.join(PROJECT_ROOT, "Publishing-PWA-to-Microsoft-Store.pptx")

# --- Theme ------------------------------------------------------------------
DARK_BG = RGBColor(0x11, 0x18, 0x2B)      # deep navy for the title slide
ACCENT = RGBColor(0x4F, 0x8C, 0xFF)       # bright blue accent
LIGHT_TEXT = RGBColor(0xF5, 0xF7, 0xFA)   # near-white
BODY_BG = RGBColor(0xFF, 0xFF, 0xFF)      # white body slides
BODY_TITLE = RGBColor(0x1B, 0x2A, 0x4A)   # dark slate for body titles
BODY_TEXT = RGBColor(0x27, 0x2E, 0x3B)    # readable dark gray
FOOTER_TEXT = RGBColor(0x8A, 0x93, 0xA5)  # muted gray footer

FOOTER = "Publishing a PWA to the Microsoft Store"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# --- Title slide ------------------------------------------------------------
TITLE_SLIDE = {
    "title": "Publishing a PWA to the Microsoft Store",
    "subtitle": "A repeatable process for shipping any Progressive Web App",
    "tagline": "Built with GitHub Copilot",
}

# === MILESTONES (append new slides here) ====================================
# Each dict = one content slide: a "title" plus a list of "bullets".
# Optional "image" key (absolute path) embeds a picture on the right side.
SLIDES = [
    {
        "title": "Why publish a PWA to the Store?",
        "bullets": [
            "Reach Windows users through the Microsoft Store with no native rewrite",
            "One web codebase powers the browser experience and the Store app",
            "Content updates flow automatically from your HTTPS host - no re-submit for most changes",
            "Installable, offline-capable, and app-like via the manifest + service worker",
        ],
    },
    {
        "title": "The end-to-end process",
        "bullets": [
            "1. Make the web app a compliant, installable PWA",
            "2. Host it on a public HTTPS URL",
            "3. Package it for Windows with PWABuilder (produces a signed MSIX)",
            "4. Submit the package in Partner Center and pass certification",
            "5. Maintain: push web updates; repackage only when the app version changes",
        ],
    },
    {
        "title": "Step 1 - PWA compliance checklist",
        "bullets": [
            "Web App Manifest with: id, name, short_name, description, start_url, scope",
            "display: standalone, plus theme_color and background_color",
            "Icons at 192 and 512 px, including a 'maskable' purpose",
            "Screenshots for narrow (mobile) and wide (desktop) form factors",
            "A registered service worker (enables offline + installability)",
            "Served over HTTPS - a hard requirement for PWAs and the Store",
        ],
    },
    {
        "title": "The three pillars PWABuilder checks",
        "bullets": [
            "Manifest = identity: a JSON file (manifest.json) describing the app - name, icons, colors, and how it launches. It makes a website installable and gives PWABuilder the metadata to build the Store package.",
            "Service Worker = offline + installable: a background script (sw.js) between the app and the network. It caches files so the app loads fast and works offline, and it enables installation. Prefer network-first: fresh when online, cached fallback offline.",
            "Security = HTTPS: the app must be served over HTTPS. Service workers only run on secure origins and the Store requires it. Hosting on GitHub Pages satisfies this automatically.",
            "In short: Manifest = identity, Service Worker = offline + installable, Security = HTTPS.",
        ],
    },
    {
        "title": "Common gotcha - service worker caching",
        "bullets": [
            "Symptom: code/asset edits don't appear after refresh (stale files served)",
            "Cause: a cache-first service worker keeps serving old cached assets",
            "Fix: use network-first (fresh when online, cache as offline fallback)",
            "Bump the cache version (CACHE_NAME) whenever cached assets change",
            "Hard-refresh (Ctrl+Shift+R) to clear the browser's own HTTP cache",
        ],
    },
    {
        "title": "Step 2 - Host on public HTTPS",
        "bullets": [
            "You need a stable, public HTTPS URL for PWABuilder to read",
            "Free options: GitHub Pages, Azure Static Web Apps, Netlify, Vercel, Cloudflare Pages",
            "GitHub Pages: push to a repo, let a GitHub Actions workflow auto-deploy on each push",
            "Use relative asset paths so the app also works under a /repo-name/ subpath",
        ],
    },
    {
        "title": "Step 3 - Package with PWABuilder",
        "bullets": [
            "Open pwabuilder.com and enter your live HTTPS URL",
            "It audits the manifest, service worker, and security, and reports gaps to fix",
            "Choose Package For Stores -> Windows to generate a signed .msixbundle",
            "Supply your Partner Center Publisher ID + display name so the package matches your account",
        ],
    },
    {
        "title": "The exact role of PWABuilder",
        "bullets": [
            "PWABuilder is the bridge between your web app and the Microsoft Store",
            "1. Audited the PWA: fetched the live URL and validated manifest, service worker, and HTTPS",
            "2. Wrapped the web app in a Windows app container (the core job): generated the native AppxManifest, embedded your Store identity (Package ID, Publisher ID), and packed icons and metadata",
            "3. Produced the deliverables: .msixbundle (main), .classic.appxbundle (compatibility), plus sideload test artifacts",
            "What it does NOT do: it does not host your app - your code still runs from GitHub Pages",
            "The generated package is a thin native shell that loads your live HTTPS site in a Windows WebView",
            "Push an update to GitHub and installed users get it automatically; PWABuilder does not run or maintain your code, and does not submit for you",
            "In one line: PWABuilder turns a live PWA URL into signed, Store-ready Windows packages - a packager, not a host or runtime",
        ],
    },
    {
        "title": "Where PWABuilder sits in the pipeline",
        "bullets": [
            "Step 1: Your code is hosted on GitHub Pages over HTTPS, giving a live URL",
            "Step 2: PWABuilder reads that live URL, audits it, and wraps it into MSIX and APPX packages with your Store identity embedded",
            "Step 3: You upload the .msixbundle and .appxbundle to Partner Center, complete the listing, and submit",
            "Step 4: The Microsoft Store distributes the app; users install a native shell that loads your live site",
            "Key idea: hosting (GitHub Pages), packaging (PWABuilder), and distribution (Partner Center / Store) are three separate responsibilities",
        ],
    },
    {
        "title": "Matching the package identity to Partner Center",
        "bullets": [
            "PWABuilder shows example placeholder values - you must overwrite them with your real Partner Center identity",
            "Package ID  =  Partner Center 'Package/Identity/Name' (the unique app identity the Store assigns)",
            "Publisher ID  =  Partner Center 'Package/Identity/Publisher' (the CN=... value)",
            "Publisher display name  =  Partner Center 'Package/Properties/PublisherDisplayName'",
            "These must match exactly or Partner Center rejects the upload",
            "Set Version to 1.0.1 or higher - the Store reserves 1.0.0.0",
        ],
    },
    {
        "title": "What is inside the PWABuilder package",
        "bullets": [
            ".msixbundle - the main Store package; upload this to Partner Center",
            ".classic.appxbundle - compatibility package for older Windows 10; also upload this to Partner Center",
            ".sideload.msix - a test copy to install locally (bypasses the Store); not uploaded",
            "install.ps1 - helper script that installs the sideload package and its test certificate for local testing",
            "utils folder - support files (test certificate and helpers) used by install.ps1",
            "readme.html - PWABuilder's own instructions for these files",
            "To submit: upload BOTH the .msixbundle and the .classic.appxbundle",
        ],
    },
    {
        "title": "Windows package formats the Store accepts",
        "bullets": [
            "MSIX family is the modern format (2018+); APPX family is the older Windows 8/10 UWP era format - same idea, newer container",
            ".msix - a single app package for one architecture (modern)",
            ".msixbundle - a bundle of several .msix packages (multiple architectures and scales); PWABuilder's main output and what you upload",
            ".msixupload - a Store submission wrapper around a .msixbundle plus debug symbols; produced by Visual Studio, not needed here",
            ".appx - the original single-app package for one architecture (older)",
            ".appxbundle - a bundle of .appx packages; PWABuilder's classic compatibility package, also uploaded",
            ".appxupload - the Visual Studio wrapper around an .appxbundle plus symbols; not needed here",
            ".xap - legacy Windows Phone 7/8 format, essentially dead; ignore it",
            "Rule of thumb: msix = new, appx = old, bundle = multiple architectures, upload = VS wrapper with symbols",
        ],
    },
    {
        "title": "Two upload questions clarified",
        "bullets": [
            "Why pick a device family if the package already declares one?",
            "Package manifest (TargetDeviceFamily) = the technical capability: which OS families the package CAN run on",
            "Partner Center checkboxes = your distribution choice: which device families you WANT to offer the app to",
            "They are separate on purpose - e.g. a package could run on Xbox but you may choose not to support it",
            "If no box is checked, the Store has no market to list the app in, so you must opt in (check Windows 10/11 Desktop)",
            "Why is .msixupload not needed for a PWA?",
            "It is a Visual Studio wrapper that bundles native debug symbols (.appxsym) for crash analysis",
            "PWABuilder does not produce one, and the Store accepts .msixbundle directly",
            "A PWA is HTML, CSS, and JavaScript with no native compiled code, so there are no symbols to ship - the format is irrelevant",
        ],
    },
    {
        "title": "Milestone - Packaging in progress",
        "bullets": [
            "Live app validated and hosted at https://vijayatmicrosoft.github.io/number-tiles-pwa/",
            "Privacy policy published at /privacy.html (required for the Store listing)",
            "Ran PWABuilder against the live URL to generate the Windows package",
            "Key step: Package ID, Publisher ID, and Publisher display name must match the reserved app in Partner Center exactly",
            "Output: a signed .msixbundle ready to upload to Partner Center",
        ],
    },
    {
        "title": "Step 4 - Submit in Partner Center",
        "bullets": [
            "Register a developer account once (individual ~$19 / company ~$99)",
            "Reserve the app name, then upload the .msixbundle under Packages",
            "Complete the listing: description, screenshots, category, age rating",
            "Provide a privacy policy URL - it is required to pass certification",
            "Submit and wait for Store certification review",
        ],
    },
    {
        "title": "Step 5 - Maintain & update",
        "bullets": [
            "Content/UI changes: just deploy to your HTTPS host - users get them automatically",
            "Bump the service worker cache version so clients pick up new assets",
            "Only repackage + resubmit when the app package version/identity changes",
            "Watch Partner Center for certification status, ratings, and health metrics",
        ],
    },
    {
        "title": "Reusable checklist",
        "bullets": [
            "[ ] Manifest complete (icons 192/512 + maskable, screenshots, colors, scope)",
            "[ ] Service worker registered and network-first for active development",
            "[ ] Deployed to a public HTTPS URL",
            "[ ] Packaged to .msixbundle via PWABuilder with correct Publisher identity",
            "[ ] Privacy policy page published and linked",
            "[ ] Submitted in Partner Center and passed certification",
        ],
    },
    {"section": True, "title": "Code Signing and Certificates"},
    {
        "title": "Why Windows requires signing",
        "bullets": [
            "Windows refuses to install an unsigned MSIX or APPX package",
            "A signature answers two questions: who made this (authenticity) and was it changed since (integrity)",
            "Think of it as a tamper-proof wax seal stamped with a unique signet ring",
        ],
    },
    {
        "title": "How a digital signature works",
        "bullets": [
            "Uses asymmetric cryptography: a private key (secret) and a public key (shared)",
            "Sign: hash the package with SHA-256, then encrypt that hash with the private key to produce the signature",
            "Verify: Windows re-hashes the package and checks it against the signature using the public key",
            "If they match, the package is unmodified and was signed by the holder of the private key",
        ],
    },
    {
        "title": "Certificates and the chain of trust",
        "bullets": [
            "A public key alone does not prove identity",
            "A certificate binds a public key to an identity and is vouched for by a Certificate Authority (CA)",
            "Windows trusts a set of root CAs and walks the chain: package to your certificate to a trusted root",
            "A self-signed certificate has no CA backing, so Windows shows an untrusted publisher warning",
        ],
    },
    {
        "title": "Why you did not need your own certificate",
        "bullets": [
            "For Store apps, Microsoft signs the package for you during certification",
            "The trust chain becomes: Microsoft Root CA to Microsoft Store CA to your app",
            "Every Windows device already trusts Microsoft's root, so the app installs cleanly",
            "This is why the Publisher ID (CN=GUID) had to match: it is your identity in Microsoft's signing system",
            "You would only buy a code-signing certificate if you distributed the app outside the Store",
        ],
    },
    {
        "title": "The test certificate in your package",
        "bullets": [
            "PWABuilder also generated a self-signed test certificate (in the utils folder, used by install.ps1)",
            "install.ps1 trusts that certificate on your machine, then installs the sideload .msix for local testing",
            "Self-signed is fine on your own PC but is not trusted by anyone else",
            "Public trust comes from Microsoft's signature, not yours",
            "Bonus: timestamping proves it was signed while the certificate was valid, so it keeps working after the cert expires",
        ],
    },
]
# === END MILESTONES =========================================================


def _add_footer(slide):
    """Add a muted footer line to a slide."""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.02), Inches(12.333), Inches(0.35))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = FOOTER
    run.font.size = Pt(10)
    run.font.color.rgb = FOOTER_TEXT
    p.alignment = PP_ALIGN.LEFT


def _fit_image(slide, image_path, left, top, max_w, max_h):
    """Add an image scaled to fit within the given box, preserving aspect ratio."""
    pic = slide.shapes.add_picture(image_path, left, top)
    scale = min(max_w / pic.width, max_h / pic.height)
    new_w = int(pic.width * scale)
    new_h = int(pic.height * scale)
    pic.width = new_w
    pic.height = new_h
    # Center within the target box.
    pic.left = int(left + (max_w - new_w) / 2)
    pic.top = int(top + (max_h - new_h) / 2)
    return pic


def build_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Dark background rectangle covering the whole slide.
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)  # 1 = rectangle
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    bg.shadow.inherit = False

    # Accent bar.
    bar = slide.shapes.add_shape(1, Inches(0.9), Inches(2.55), Inches(2.4), Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False

    # Title text.
    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.8), Inches(11.5), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = TITLE_SLIDE["title"]
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = LIGHT_TEXT

    # Subtitle.
    sub_box = slide.shapes.add_textbox(Inches(0.9), Inches(4.5), Inches(11.5), Inches(0.9))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = TITLE_SLIDE["subtitle"]
    r.font.size = Pt(22)
    r.font.color.rgb = RGBColor(0xC5, 0xD1, 0xE8)

    # Tagline.
    tag_box = slide.shapes.add_textbox(Inches(0.9), Inches(5.4), Inches(11.5), Inches(0.6))
    tf = tag_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = TITLE_SLIDE["tagline"]
    r.font.size = Pt(14)
    r.font.italic = True
    r.font.color.rgb = ACCENT


def build_section_slide(prs, data):
    """Render a dark section-divider slide with a centered title and accent bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Dark background rectangle covering the whole slide.
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)  # 1 = rectangle
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    bg.shadow.inherit = False

    # Centered accent bar above the title.
    bar_w = Inches(2.4)
    bar = slide.shapes.add_shape(
        1, int((SLIDE_W - bar_w) / 2), Inches(2.85), bar_w, Inches(0.12)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False

    # Centered section title in large light text.
    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(3.2), Inches(11.5), Inches(1.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = data["title"]
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = LIGHT_TEXT


def build_content_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # White background.
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BODY_BG
    bg.line.fill.background()
    bg.shadow.inherit = False

    has_image = bool(data.get("image")) and os.path.exists(data["image"])

    # Title.
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.1), Inches(1.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = data["title"]
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = BODY_TITLE

    # Accent underline.
    ul = slide.shapes.add_shape(1, Inches(0.62), Inches(1.45), Inches(2.0), Inches(0.06))
    ul.fill.solid()
    ul.fill.fore_color.rgb = ACCENT
    ul.line.fill.background()
    ul.shadow.inherit = False

    # Body text width depends on whether an image sits on the right.
    body_w = Inches(7.4) if has_image else Inches(12.1)
    body_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), body_w, Inches(4.9))
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, bullet in enumerate(data["bullets"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = "\u2022  " + bullet
        r.font.size = Pt(18)
        r.font.color.rgb = BODY_TEXT
        p.space_after = Pt(12)

    # Optional image on the right.
    if has_image:
        _fit_image(
            slide,
            data["image"],
            left=Inches(8.3),
            top=Inches(1.8),
            max_w=Inches(4.4),
            max_h=Inches(4.9),
        )

    _add_footer(slide)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs)
    for data in SLIDES:
        if data.get("section"):
            build_section_slide(prs, data)
        else:
            build_content_slide(prs, data)

    prs.save(OUTPUT_PATH)

    total_slides = len(prs.slides._sldIdLst)
    print("SUCCESS: deck generated.")
    print("File: " + os.path.abspath(OUTPUT_PATH))
    print("Total slides: " + str(total_slides))


if __name__ == "__main__":
    main()
